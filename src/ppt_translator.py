from googletrans import Translator
from pptx import Presentation

title = input('name of ppt (ex: Policies 2021): ')

t = Translator()

prs = Presentation(title+'.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if len(paragraph.text)==0:
                    continue

                print(paragraph.text)
                
                # Store font of original paragraph  (eg. Chinese)
                init_font = paragraph.runs[0].font  # 'runs' is a tuple, where runs[0] is the paragraph. runs[1] doesn't exist.
                FName = init_font.name
                FSize = init_font.size
                FBold = init_font.bold

                # Translate
                cur_text = paragraph.text
                new_text = t.translate(cur_text, dest='en').text
                paragraph.text = new_text

                # Apply font to translated paragraph  (eg. English)
                new_par = paragraph.runs[0]
                new_par.font.name = FName
                new_par.font.size = FSize
                new_par.font.bold = FBold

prs.save('Translated_'+title+'.pptx')
